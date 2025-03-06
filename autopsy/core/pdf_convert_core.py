import os
import fitz  # PyMuPDF
from PIL import Image
import io

# If available, you can use pdf2docx for PDF to DOCX conversion.
try:
    from pdf2docx import Converter
except ImportError:
    Converter = None

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None

def convert_pdf(input_path, output_folder, conversion_type, progress_callback=None, image_format="png"):
    """
    Convert a PDF file to a different format.
    
    conversion_type: one of "docx", "ppt", "images".
      - "docx": convert entire PDF to DOCX using pdf2docx.
      - "ppt": create a PPTX where each slide is an image of a PDF page.
      - "images": export each page as an individual image file.
    
    image_format (for images mode): one of "jpg", "png", "bmp" (default "png")
    
    Returns:
        list: A list of file paths for the generated files.
    """
    if not os.path.isfile(input_path):
        raise FileNotFoundError(f"Input file not found: {input_path}")
    if not os.path.isdir(output_folder):
        os.makedirs(output_folder)
    
    if conversion_type.lower() == "docx":
        if Converter is None:
            raise ImportError("pdf2docx is not installed. Please install it to convert PDF to DOCX.")
        output_file = os.path.join(output_folder, os.path.splitext(os.path.basename(input_path))[0] + ".docx")
        cv = Converter(input_path)
        cv.convert(output_file, start=0, end=None)
        cv.close()
        if progress_callback:
            progress_callback(100)
        return [output_file]

    elif conversion_type.lower() == "ppt":
        if Presentation is None:
            raise ImportError("python-pptx is not installed. Please install it to convert PDF to PPT.")
        prs = Presentation()
        doc = fitz.open(input_path)
        total_pages = len(doc)
        output_file = os.path.join(output_folder, os.path.splitext(os.path.basename(input_path))[0] + ".pptx")
        for i in range(total_pages):
            pdf_page = doc[i]
            zoom = 2  # Increase zoom for better quality
            mat = fitz.Matrix(zoom, zoom)
            pix = pdf_page.get_pixmap(matrix=mat)
            temp_image = os.path.join(output_folder, f"temp_page_{i+1}.png")
            pix.save(temp_image)
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
            slide.shapes.add_picture(temp_image, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            os.remove(temp_image)
            if progress_callback:
                progress_callback(int(((i + 1)/total_pages)*100))
        prs.save(output_file)
        return [output_file]

    elif conversion_type.lower() == "images":
        doc = fitz.open(input_path)
        total_pages = len(doc)
        output_files = []
        # Ensure the image format is in lower-case.
        ext = image_format.lower()
        for i in range(total_pages):
            pdf_page = doc[i]
            pix = pdf_page.get_pixmap()
            output_filename = os.path.join(
                output_folder,
                f"{os.path.splitext(os.path.basename(input_path))[0]}_page_{i+1}.{ext}"
            )
            if ext == "png":
                pix.save(output_filename)
            else:
                # For JPG or BMP, convert via PIL.
                img_data = pix.tobytes("ppm")
                from PIL import Image
                import io
                pil_img = Image.open(io.BytesIO(img_data))
                # Use "JPEG" for jpg images.
                format_str = "JPEG" if ext == "jpg" else ext.upper()
                pil_img.save(output_filename, format=format_str)
            output_files.append(output_filename)
            if progress_callback:
                progress_callback(int(((i + 1)/total_pages)*100))
        return output_files

    else:
        raise ValueError(f"Unknown conversion type: {conversion_type}")
